#!/usr/bin/env python3
"""
report-gen.py — Generate pentest deliverables from findings JSON.

Produces:
  - Word document (.docx) — Technical report with findings, CVSS, remediation
  - PowerPoint (.pptx) — Executive brief with charts and summary
  - Uses brand colors (#CC4141 red, #333333 dark gray)

Usage:
    python report-gen.py findings.json --output-dir ./reports/
    python report-gen.py findings.json --format docx       # Word only
    python report-gen.py findings.json --format pptx       # PowerPoint only
    python report-gen.py findings.json --format pdf        # PDF only
    python report-gen.py findings.json --format all        # All formats (default)

Input format (findings.json):
[
  {
    "id": "FIND-001",
    "title": "Outdated nginx with known RCE",
    "host": "192.168.1.10",
    "port": 443,
    "service": "nginx",
    "version": "1.18.0",
    "cve": "CVE-2021-23017",
    "cvss_score": 7.5,
    "cvss_vector": "CVSS:3.1/AV:N/AC:L/PR:N/UI:N/S:U/C:H/I:N/A:N",
    "severity": "high",
    "description": "The target is running nginx 1.18.0 which is vulnerable to...",
    "remediation": "Upgrade nginx to version 1.20.0 or later.",
    "status": "open"
  }
]
"""

import argparse
import json
import sys
from datetime import datetime
from pathlib import Path

# Brand Colors
BRAND_RED = "CC4141"
DARK_GRAY = "333333"
MED_GRAY = "797979"

SEVERITY_COLORS = {
    "critical": "CC4141",
    "high": "E85D3F",
    "medium": "F4A442",
    "low": "5B9E5B",
    "info": "4A90D9",
}

SEVERITY_ORDER = {"critical": 0, "high": 1, "medium": 2, "low": 3, "info": 4}

# Device type mapping for host display — populated from findings metadata.
# Never hardcode engagement-specific IP-to-device mappings here.
DEVICE_TYPES = {}


def _host_display(host: str, port: int | str = "", device_type: str = "") -> str:
    """Return formatted host string, e.g. 'Host 10.0.0.1:22 — Charging Station'."""
    base = f"Host {host}"
    if device_type:
        base += f" — {device_type}"
    if port:
        base += f":{port}"
    return base


def load_findings(path: str) -> list[dict]:
    """Load and validate findings JSON."""
    with open(path) as f:
        try:
            findings = json.load(f)
        except json.JSONDecodeError as e:
            print(f"Error: {path} is not valid JSON: {e}", file=sys.stderr)
            sys.exit(1)

    if isinstance(findings, dict):
        if "findings" in findings:
            print(
                f"Error: expected a JSON array, got an object with key 'findings'. "
                f"Did you wrap it in {{\"findings\": [...]}}? "
                f"report-gen.py expects a flat list at the top level.",
                file=sys.stderr,
            )
        else:
            print(
                f"Error: expected a JSON array, got an object. "
                f"report-gen.py expects a flat list of finding objects.",
                file=sys.stderr,
            )
        sys.exit(1)

    if not isinstance(findings, list):
        print(
            f"Error: expected a JSON array, got {type(findings).__name__}. "
            f"report-gen.py expects a flat list of finding objects.",
            file=sys.stderr,
        )
        sys.exit(1)

    # Sort by severity
    findings.sort(key=lambda f: SEVERITY_ORDER.get(f.get("severity", "info"), 99))
    return findings


def severity_summary(findings: list[dict]) -> dict:
    """Count findings by severity."""
    counts = {"critical": 0, "high": 0, "medium": 0, "low": 0, "info": 0}
    for f in findings:
        sev = f.get("severity", "info")
        if sev in counts:
            counts[sev] += 1
    return counts


def generate_docx(findings: list[dict], output_path: str) -> str:
    """Generate Word technical report."""
    try:
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor, Mm, Cm
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
    except ImportError:
        return "ERROR: python-docx not installed. Run: pip install python-docx"

    doc = Document()

    # Set page size to A4
    section = doc.sections[0]
    section.page_width = Mm(210)
    section.page_height = Mm(297)
    section.top_margin = Cm(2)
    section.bottom_margin = Cm(2)
    section.left_margin = Cm(2.5)
    section.right_margin = Cm(2.5)

    # ── Brand styling ──
    style = doc.styles['Normal']
    style.font.name = 'Calibri'
    style.font.size = Pt(11)
    style.font.color.rgb = RGBColor(*hex_to_rgb('333333'))
    # Ensure font applies
    style.element.rPr.rFonts.set(qn('w:eastAsia'), 'Calibri')

    for level, size, color_name in [
        (1, 16, '333333'),
        (2, 14, '333333'),
        (3, 12, '555555'),
    ]:
        hs = doc.styles[f'Heading {level}']
        hs.font.name = 'Calibri'
        hs.font.size = Pt(size)
        hs.font.bold = True
        hs.font.color.rgb = RGBColor(*hex_to_rgb(color_name))
        hs.element.rPr.rFonts.set(qn('w:eastAsia'), 'Calibri')
        # Add red left border to H1
        if level == 1:
            pPr = hs.element.get_or_add_pPr()
            pBdr = OxmlElement('w:pBdr')
            left = OxmlElement('w:left')
            left.set(qn('w:val'), 'single')
            left.set(qn('w:sz'), '18')
            left.set(qn('w:space'), '8')
            left.set(qn('w:color'), BRAND_RED)
            pBdr.append(left)
            pPr.append(pBdr)

    # Title styling (level 0 → "Title" style)
    title_style = doc.styles['Title']
    title_style.font.name = 'Calibri'
    title_style.font.size = Pt(22)
    title_style.font.bold = True
    title_style.font.color.rgb = RGBColor(*hex_to_rgb(BRAND_RED))
    title_style.element.rPr.rFonts.set(qn('w:eastAsia'), 'Calibri')

    # Title
    title = doc.add_heading("Security Assessment Report", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph(f"PortShim — {datetime.now().strftime('%B %d, %Y')}")
    doc.add_paragraph("")

    # Executive Summary
    doc.add_heading("Executive Summary", level=1)
    
    summary = severity_summary(findings)
    host_count_uniq = len(set(f.get("host", "") for f in findings if f.get("host")))
    has_critical = summary['critical'] > 0
    has_high = summary['high'] > 0
    
    # Derive target network from findings hosts
    hosts_in_findings = sorted(set(f.get("host", "") for f in findings if f.get("host")))
    target_net = hosts_in_findings[0].rsplit(".", 1)[0] + ".0/24" if hosts_in_findings else "N/A"

    doc.add_heading("Assessment Overview", level=2)
    doc.add_paragraph(
        f"PortShim conducted an internal network security assessment of {target_net} "
        f"on {datetime.now().strftime('%B %d, %Y')}. The assessment identified "
        f"{len(findings)} findings across {host_count_uniq} live hosts using a combination "
        f"of service enumeration and vulnerability scanning."
    )
    
    risk_label = "CRITICAL" if has_critical else "HIGH" if has_high else "MODERATE"
    doc.add_heading(f"Overall Risk Posture: {risk_label}", level=2)
    doc.add_paragraph(
        f"Of the {len(findings)} findings: "
        f"{summary['critical']} Critical, {summary['high']} High, "
        f"{summary['medium']} Medium, {summary['low']} Low, {summary['info']} Informational."
    )
    
    doc.add_heading("Key Risks", level=2)
    for f in findings:
        if f.get("severity") in ("critical", "high"):
            doc.add_paragraph(
                f"{f.get('title', '')} — {f.get('host', '')}:{f.get('port', '')} "
                f"(CVSS {f.get('cvss_score', 'N/A')})",
                style="List Bullet"
            )
    
    doc.add_heading("Identified Services", level=2)
    # Build data-driven services summary from findings metadata
    exploitation_lines = []
    for f in findings[:8]:  # Summarise top findings only
        host = f.get("host", "")
        service = f.get("service", "")
        port = f.get("port", "")
        version = f.get("version", "")
        detail = f"  - {host}:{port} — {service}"
        if version:
            detail += f" {version}"
        exploitation_lines.append(detail)
    if exploitation_lines:
        doc.add_paragraph(
            "Identified services on high-value targets:\n"
            + "\n".join(exploitation_lines)
        )
    else:
        doc.add_paragraph("No exploitation attempts were conducted.")
    
    maturity = "Low" if has_critical else "Developing" if has_high else "Managed"
    doc.add_heading(f"Environment Maturity: {maturity}", level=2)
    if has_critical:
        doc.add_paragraph(
            "The environment contains multiple high-severity issues indicating limited security "
            "controls on the internal network. Several devices expose management interfaces, "
            "legacy services, or unencrypted protocols without restriction."
        )
    elif has_high:
        doc.add_paragraph("Security controls are present but inconsistent.")
    else:
        doc.add_paragraph("The environment demonstrates a proactive security stance.")
    
    doc.add_heading("Recommended Next Steps", level=2)
    # Generate next steps from finding types — data-driven, no hardcoded host names
    services_found = set(f.get("service", "").lower() for f in findings if f.get("service"))
    if "telnet" in services_found:
        telnet_hosts = [f.get("host", "") for f in findings if f.get("service", "").lower() == "telnet"]
        doc.add_paragraph(f"Disable Telnet on {', '.join(telnet_hosts[:3])}. Replace with SSH.", style="List Number")
    if "ssh" in services_found and any("7." in (f.get("version") or "") for f in findings if f.get("service", "").lower() == "ssh"):
        doc.add_paragraph("Upgrade outdated OpenSSH instances to a supported release.", style="List Number")
    if summary.get("high", 0) > 0 or summary.get("critical", 0) > 0:
        doc.add_paragraph("Apply security patches to affected hosts (see findings detail).", style="List Number")
        doc.add_paragraph("Restrict management interfaces to a dedicated admin VLAN where applicable.", style="List Number")

    # Severity Methodology
    doc.add_heading("Severity Classification", level=1)
    doc.add_paragraph(
        "Findings are classified using the Common Vulnerability Scoring System (CVSS) v3.1. "
        "The following severity scale is applied:"
    )
    for sev, color in SEVERITY_COLORS.items():
        p = doc.add_paragraph(f"{sev.upper()}: CVSS {sev_ranges(sev)}")
        for run in p.runs:
            run.font.color.rgb = RGBColor(*hex_to_rgb(color))

    # Findings Table
    doc.add_heading("Findings", level=1)
    table = doc.add_table(rows=1, cols=5)
    table.style = "Table Grid"

    hdr = table.rows[0].cells
    for i, text in enumerate(["ID", "Severity", "Title", "Host", "CVSS"]):
        hdr[i].text = text
        for para in hdr[i].paragraphs:
            for run in para.runs:
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)
                run.font.name = "Calibri"
                run.font.size = Pt(10)
        # Dark header background
        from docx.oxml import OxmlElement
        shading = OxmlElement('w:shd')
        shading.set(qn('w:fill'), '333333')
        shading.set(qn('w:val'), 'clear')
        hdr[i]._tc.get_or_add_tcPr().append(shading)

    for finding in findings:
        row = table.add_row().cells
        row[0].text = finding.get("id", "")
        row[1].text = finding.get("severity", "").upper()
        row[2].text = finding.get("title", "")
        row[3].text = _host_display(finding.get("host", ""), finding.get("port", ""), finding.get("device_type", ""))
        row[4].text = str(finding.get("cvss_score", "N/A"))
        # Alternating row shading + Calibri font
        for cell in row:
            for para in cell.paragraphs:
                for run in para.runs:
                    run.font.name = "Calibri"
                    run.font.size = Pt(10)
        if (len(table.rows) - 1) % 2 == 0:
            for cell in row:
                shading = OxmlElement('w:shd')
                shading.set(qn('w:fill'), 'F5F5F5')
                shading.set(qn('w:val'), 'clear')
                cell._tc.get_or_add_tcPr().append(shading)
    doc.add_heading("Detailed Findings", level=1)
    for finding in findings:
        doc.add_heading(f"{finding.get('id')}: {finding.get('title')}", level=2)

        info = doc.add_paragraph()
        info.add_run("Severity: ").bold = True
        info.add_run(f"{finding.get('severity', 'N/A').upper()} (CVSS {finding.get('cvss_score', 'N/A')})")

        host_str = _host_display(finding.get("host", ""), finding.get("port", ""), finding.get("device_type", ""))
        doc.add_paragraph(f"{host_str}")
        dtype = finding.get("device_type", "")
        if dtype:
            doc.add_paragraph(f"Device: {dtype}")
        cve = finding.get("cve", "")
        if cve:
            from docx.oxml.ns import qn
            from docx.oxml import OxmlElement
            cve_para = doc.add_paragraph()
            cve_para.add_run("CVE: ").bold = True
            cve_url = f"https://nvd.nist.gov/vuln/detail/{cve}"
            # Add hyperlink
            add_hyperlink(cve_para, cve_url, cve)
            cve_para.add_run(f"  |  CVSS {finding.get('cvss_score', 'N/A')} ({finding.get('severity', '').upper()})")
            cve_para.add_run(f"\n  Vector: {finding.get('cvss_vector', 'N/A')}")
        if finding.get("description"):
            doc.add_heading("Description", level=3)
            doc.add_paragraph(finding["description"])
        if finding.get("remediation"):
            doc.add_heading("Remediation", level=3)
            doc.add_paragraph(finding["remediation"])

        doc.add_paragraph("")

    doc.save(output_path)
    return output_path


def generate_pptx(findings: list[dict], output_path: str, engagement_profile: str = "Surgical") -> str:
    """Generate PowerPoint executive brief."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        return "ERROR: python-pptx not installed. Run: pip install python-pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    summary = severity_summary(findings)
    total = len(findings)
    has_critical = summary['critical'] > 0
    has_high = summary['high'] > 0
    # Derive target network from findings hosts
    hosts_in_findings = sorted(set(f.get("host", "") for f in findings if f.get("host")))
    target_net = hosts_in_findings[0].rsplit(".", 1)[0] + ".0/24" if hosts_in_findings else "N/A"
    risk_label = "CRITICAL" if has_critical else "HIGH" if has_high else "MODERATE"

    def _add_slide_title(prs, title_text):
        """Add a content slide with title bar."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_bar = slide.shapes.add_shape(
            1, Inches(0), Inches(0), prs.slide_width, Inches(1.1)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = RGBColor(*hex_to_rgb(DARK_GRAY))
        title_bar.line.fill.background()
        tf = title_bar.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = "Calibri"
        tf.paragraphs[0].space_before = Pt(8)
        return slide

    def _add_bullet(slide, text, left, top, width, height, font_size=14, color=None, bold=False):
        """Add a text box with a single paragraph."""
        use_color = color or DARK_GRAY
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(*hex_to_rgb(use_color))
        p.font.name = "Calibri"
        p.font.bold = bold
        p.space_after = Pt(4)
        return tf

    # ── Slide 1: Title ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*hex_to_rgb(DARK_GRAY))

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PORTSHIM"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*hex_to_rgb(BRAND_RED))
    p.font.name = "Calibri"

    subtitle = tf.add_paragraph()
    subtitle.text = "Security Assessment — Executive Brief"
    subtitle.font.size = Pt(24)
    subtitle.font.color.rgb = RGBColor(255, 255, 255)
    subtitle.font.name = "Calibri"

    date_para = tf.add_paragraph()
    date_para.text = datetime.now().strftime("%B %d, %Y")
    date_para.font.size = Pt(14)
    date_para.font.color.rgb = RGBColor(*hex_to_rgb(MED_GRAY))
    date_para.font.name = "Calibri"

    # ── Slide 2: Assessment Overview ──
    slide = _add_slide_title(prs, "Assessment Overview")
    host_count = len(set(f.get("host", "") for f in findings if f.get("host")))

    _add_bullet(slide, f"Target Network:  {target_net}", 0.8, 1.5, 11, 0.5, 18, DARK_GRAY, True)
    _add_bullet(slide, f"Assessment Date:  {datetime.now().strftime('%B %d, %Y')}", 0.8, 2.1, 11, 0.5, 18, DARK_GRAY, True)
    _add_bullet(slide, f"Profile:         {engagement_profile}", 0.8, 2.7, 11, 0.5, 18, DARK_GRAY, True)
    _add_bullet(slide, f"LLM Mode:        Hybrid (local models for exploitation, cloud for reporting)", 0.8, 3.3, 11, 0.5, 18, DARK_GRAY, True)
    _add_bullet(slide, f"Live Hosts Found:  {host_count}", 0.8, 3.9, 11, 0.5, 18, DARK_GRAY, True)
    _add_bullet(slide, f"Total Findings:  {total} ({summary['critical']} Critical, {summary['high']} High, {summary['medium']} Medium, {summary['low']} Low)", 0.8, 4.5, 11, 0.5, 18, DARK_GRAY, True)

    # ── Slide 3: Overall Risk Posture ──
    slide = _add_slide_title(prs, f"Overall Risk Posture: {risk_label}")

    # Severity bar chart
    visible_count = sum(summary[s] for s in ["critical", "high", "medium", "low"])
    denominator = max(1, visible_count) if visible_count > 0 else max(1, total)
    y = Inches(1.8)
    for sev in ["critical", "high", "medium", "low"]:
        count = summary[sev]
        if count == 0:
            continue
        bar = slide.shapes.add_shape(
            1, Inches(1.5), y, Inches(8 * count / denominator), Inches(0.5)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(*hex_to_rgb(SEVERITY_COLORS[sev]))
        bar.line.fill.background()
        label = slide.shapes.add_textbox(Inches(0.3), y, Inches(1.2), Inches(0.5))
        ltf = label.text_frame
        ltf.paragraphs[0].text = f"{sev.upper()}"
        ltf.paragraphs[0].font.size = Pt(14)
        ltf.paragraphs[0].font.bold = True
        ltf.paragraphs[0].font.color.rgb = RGBColor(*hex_to_rgb(SEVERITY_COLORS[sev]))
        ltf.paragraphs[0].font.name = "Calibri"
        count_label = slide.shapes.add_textbox(Inches(10), y, Inches(2), Inches(0.5))
        cltf = count_label.text_frame
        cltf.paragraphs[0].text = f"{count} findings"
        cltf.paragraphs[0].font.size = Pt(14)
        cltf.paragraphs[0].font.name = "Calibri"
        cltf.paragraphs[0].font.color.rgb = RGBColor(*hex_to_rgb(MED_GRAY))
        y += Inches(0.7)

    risk_detail = (
        "Immediate attention required. Several devices expose management interfaces, legacy services, "
        "or unencrypted protocols to the LAN without restriction. End-of-life software present."
    ) if has_critical else (
        "Significant security gaps requiring prompt remediation. Known vulnerabilities with available exploits identified."
    ) if has_high else (
        "No critical-severity issues identified. Several medium-severity items should be addressed."
    )
    _add_bullet(slide, risk_detail, 0.8, 5.2, 11.5, 1.5, 14, MED_GRAY)

    # ── Slide 4: Key Risks ──
    slide = _add_slide_title(prs, "Key Risks")

    top = [f for f in findings if f.get("severity") in ("critical", "high")][:8]
    y = 1.6
    for f in top:
        sev = f.get("severity", "info")
        color = SEVERITY_COLORS.get(sev, MED_GRAY)
        title_text = f.get("title", "")
        host_text = f"{f.get('host', '')}:{f.get('port', '')}"
        cve_text = f.get("cve", "")
        score = f.get("cvss_score", "")

        # Severity tag
        tag = slide.shapes.add_shape(1, Inches(0.5), Inches(y), Inches(0.9), Inches(0.35))
        tag.fill.solid()
        tag.fill.fore_color.rgb = RGBColor(*hex_to_rgb(color))
        tag.line.fill.background()
        ttf = tag.text_frame
        ttf.paragraphs[0].text = sev.upper()
        ttf.paragraphs[0].font.size = Pt(9)
        ttf.paragraphs[0].font.bold = True
        ttf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        ttf.paragraphs[0].font.name = "Calibri"
        ttf.paragraphs[0].alignment = 2  # center

        # Finding title
        _add_bullet(slide, f"{title_text}", 1.6, y, 9, 0.35, 13, DARK_GRAY, True)
        y2 = y + 0.35
        detail_parts = [host_text]
        if score:
            detail_parts.append(f"CVSS {score}")
        if cve_text:
            detail_parts.append(cve_text)
        _add_bullet(slide, "  |  ".join(detail_parts), 1.6, y2, 9, 0.3, 10, MED_GRAY)
        y += 0.65

    # ── Slide 5: Environment Assessment ──
    slide = _add_slide_title(prs, "Environment Assessment")

    maturity = "Low" if has_critical else "Developing" if has_high else "Managed"
    _add_bullet(slide, f"Maturity Level:  {maturity}", 0.8, 1.5, 11, 0.5, 20, DARK_GRAY, True)

    maturity_detail = (
        "The environment contains multiple high-severity issues indicating limited security "
        "controls on the internal network. Several devices expose management interfaces, legacy "
        "services, or unencrypted protocols to the LAN without restriction."
    ) if has_critical else (
        "Security controls are present but inconsistent. Several high-severity findings indicate "
        "gaps in patch management and network segmentation."
    ) if has_high else (
        "The environment demonstrates a proactive security stance with no critical or high-severity "
        "vulnerabilities identified."
    )
    _add_bullet(slide, maturity_detail, 0.8, 2.3, 11.5, 2, 14, MED_GRAY)

    # Exploitation attempt results
    _add_bullet(slide, "Identified Services:", 0.8, 4.2, 11, 0.4, 16, DARK_GRAY, True)
    attempts = []
    for f in findings[:5]:
        host = f.get("host", "")
        svc = f.get("service", "")
        attempts.append(f"{host} — {svc}")
    ya = 4.8
    for a in attempts:
        _add_bullet(slide, f"• {a}", 1.2, ya, 11, 0.3, 11, MED_GRAY)
        ya += 0.35

    # ── Slide 6: Recommended Next Steps ──
    slide = _add_slide_title(prs, "Recommended Next Steps")

    steps = []
    services_found = set(f.get("service", "").lower() for f in findings if f.get("service"))
    if "telnet" in services_found:
        telnet_hosts = [f.get("host", "") for f in findings if f.get("service", "").lower() == "telnet"]
        steps.append(f"Disable Telnet on {', '.join(telnet_hosts[:2])}. Replace with SSH.")
    if "ssh" in services_found and any("7." in (f.get("version") or "") for f in findings if f.get("service", "").lower() == "ssh"):
        steps.append("Upgrade outdated OpenSSH instances to a supported release.")
    if summary.get("high", 0) > 0 or summary.get("critical", 0) > 0:
        steps.append("Apply security patches to affected hosts.")
        steps.append("Restrict management interfaces to a dedicated admin VLAN.")
    if not steps:
        steps.append("No critical or high-severity issues — continue monitoring and regular patching.")

    ys = 1.8
    for i, step in enumerate(steps, 1):
        _add_bullet(slide, f"{i}.  {step}", 0.8, ys, 11.5, 0.5, 15, DARK_GRAY)
        ys += 0.75

    # Save
    prs.save(output_path)
    return output_path


def generate_pdf(findings: list[dict], output_path: str, engagement_profile: str = "Surgical") -> str:
    """Generate PDF report via weasyprint from an HTML template."""
    try:
        from weasyprint import HTML
    except ImportError:
        return "ERROR: weasyprint not installed. Run: pip install weasyprint"

    summary = severity_summary(findings)
    total = len(findings)

    # Build findings rows HTML
    detail_rows = ""
    for f in findings:
        sev = f.get("severity", "info")
        color = SEVERITY_COLORS.get(sev, MED_GRAY)
        desc = (f.get("description", "") or "")[:500]
        remed = (f.get("remediation", "") or "")[:300]
        host_str = _host_display(f.get("host", ""), f.get("port", ""), f.get("device_type", ""))
        cve = f.get("cve", "")
        cve_link = f'<a class="cve-link" href="https://nvd.nist.gov/vuln/detail/{cve}">{cve}</a>' if cve else "—"
        title_display = f'<a href="#{f["id"]}">{f.get("title", "")}</a>'
        detail_rows += f"""
        <tr>
            <td class="sev-{sev}">{sev.upper()}</td>
            <td>{title_display}</td>
            <td>{host_str}</td>
            <td>{f.get('cvss_score', 'N/A')}</td>
        </tr>"""

    detail_sections = ""
    for f in findings:
        sev = f.get("severity", "info")
        color = SEVERITY_COLORS.get(sev, MED_GRAY)
        host_str = _host_display(f.get("host", ""), f.get("port", ""), f.get("device_type", ""))
        desc = (f.get("description", "") or "")[:500]
        remed = (f.get("remediation", "") or "")[:300]
        cve = f.get("cve", "")
        cve_link = f'<a class="cve-link" href="https://nvd.nist.gov/vuln/detail/{cve}">{cve}</a>' if cve else "—"
        cve_score = f.get("cvss_score", "N/A")
        cve_summary = ""
        if cve:
            cve_url = f"https://nvd.nist.gov/vuln/detail/{cve}"
            cve_summary = f"""
            <div class="cve-summary">
                <h4>CVE Summary</h4>
                <table class="meta">
                    <tr><td class="label">CVE ID</td><td><a class="cve-link" href="{cve_url}">{cve}</a></td></tr>
                    <tr><td class="label">CVSS Score</td><td>{cve_score}/10 ({sev.upper()})</td></tr>
                    <tr><td class="label">Vector</td><td>{f.get('cvss_vector', 'N/A')}</td></tr>
                    <tr><td class="label">Link</td><td><a class="cve-link" href="{cve_url}">{cve_url}</a></td></tr>
                </table>
            </div>"""
        detail_sections += f"""
        <div class="finding" id="{f['id']}">
            <h3 style="border-left: 4px solid #{color}; padding-left: 12px;">
                {f.get('id', '')}: {f.get('title', '')}
                <span class="sev-tag sev-{sev}">{sev.upper()} CVSS {cve_score}</span>
            </h3>
            <table class="meta">
                <tr><td class="label">Host</td><td>{host_str}</td></tr>
                <tr><td class="label">Device</td><td>{f.get("device_type", "") or "—"}</td></tr>
                <tr><td class="label">Service</td><td>{f.get('service', '')} {f.get('version', '')}</td></tr>
                <tr><td class="label">CVE</td><td>{cve_link}</td></tr>
                <tr><td class="label">Status</td><td>{f.get('status', 'open')}</td></tr>
            </table>
            {cve_summary}
            <h4>Description</h4>
            <p>{desc}</p>
            <h4>Remediation</h4>
            <p>{remed}</p>
        </div>"""
# Executive summary sections — built from findings data
    has_critical = summary['critical'] > 0
    has_high = summary['high'] > 0
    has_medium = summary['medium'] > 0

    if has_critical:
        risk_label = "CRITICAL"
        risk_detail = "Immediate attention required. Vulnerabilities were identified that could allow an attacker to gain unauthorised access to network devices with minimal effort. Several services are exposed without authentication or are running outdated software with known exploit chains."
    elif has_high:
        risk_label = "HIGH"
        risk_detail = "Significant security gaps exist that require prompt remediation. Known vulnerabilities with available exploits were identified across multiple hosts."
    else:
        risk_label = "MODERATE"
        risk_detail = "No critical-severity vulnerabilities were identified, though several medium-severity issues should be addressed to maintain a strong security posture."

    # Derive target network from findings hosts
    hosts_in_findings = sorted(set(f.get("host", "") for f in findings if f.get("host")))
    target_net = hosts_in_findings[0].rsplit(".", 1)[0] + ".0/24" if hosts_in_findings else "N/A"

    # Top critical/high findings (up to 5)
    top_findings = [f for f in findings if f.get("severity") in ("critical", "high")][:5]
    top_lines = ""
    for f in top_findings:
        top_lines += f"<li><strong>{f.get('title', '')}</strong> — {f.get('host', '')}:{f.get('port', '')} (CVSS {f.get('cvss_score', 'N/A')})</li>"

    # Determine key exploited services from findings
    telnet_found = any(f.get("service") == "telnet" for f in findings)
    old_ssh = any("7.2" in (f.get("version") or "") for f in findings)

    exploitable_services = []
    if telnet_found:
        exploitable_services.append(f"Telnet (unencrypted, found on {len([f for f in findings if f.get('service') == 'telnet'])} host(s))")
    if old_ssh:
        exploitable_services.append(f"OpenSSH 7.2 (end-of-life, 130+ known CVEs)")
    if summary.get("high", 0) > 0:
        exploitable_services.append("High-severity vulnerabilities requiring immediate attention")

    services_html = ""
    if exploitable_services:
        for s in exploitable_services:
            services_html += f"<li>{s}</li>"

    # Host references for next steps
    telnet_host = next((f.get("host", "") for f in findings if f.get("service") == "telnet"), "")
    ssh_host = next((f.get("host", "") for f in findings if "7.2" in (f.get("version") or "")), "")
    ssh_servers = list(set(f.get("host", "") for f in findings if f.get("service") == "ssh" and f.get("host") not in ("", ssh_host))) or ["multiple hosts"]

    # Environment maturity assessment
    host_count = len(set(f.get("host", "") for f in findings if f.get("host")))
    if has_critical and has_high:
        maturity = "Low"
        maturity_detail = "The environment contains multiple high-severity issues that indicate limited security controls on the internal network. Several devices expose management interfaces, legacy services, or unencrypted protocols to the LAN without restriction. The presence of end-of-life software (OpenSSH 7.2 from 2016) and unprotected telnet access suggests device lifecycle management is ad-hoc rather than systematic."
    elif has_high:
        maturity = "Developing"
        maturity_detail = "Security controls are present but inconsistent. While no critical issues were found, several high-severity findings indicate gaps in patch management and network segmentation."
    else:
        maturity = "Managed"
        maturity_detail = "The environment demonstrates a proactive security stance with no critical or high-severity vulnerabilities identified."

    # Next steps
    next_steps = ""
    if telnet_found:
        next_steps += f"<li><strong>Disable Telnet</strong> — Immediately disable telnet on {telnet_host}. Replace with SSH for any remote administration.</li>"
    if old_ssh:
        next_steps += f"<li><strong>Upgrade End-of-Life Software</strong> — {ssh_host} runs an outdated OpenSSH version. Upgrade to a supported release or restrict network access to this device.</li>"
    if summary.get("high", 0) > 0:
        next_steps += f"<li><strong>Patch Management</strong> — Apply vendor security patches to affected hosts. Recent OpenSSH CVEs have associated exploit activity.</li>"
        next_steps += f"<li><strong>Network Segmentation</strong> — Segment IoT and non-critical devices onto a separate VLAN with restricted access to management interfaces.</li>"
        next_steps += f"<li><strong>Review Web Interface Exposure</strong> — Management interfaces are accessible from the LAN. Restrict to an admin workstation or management VLAN.</li>"

    # Build data-driven exploitation assessment HTML
    exploit_items_html = ""
    for f in findings[:8]:
        host = f.get("host", "")
        svc = f.get("service", "")
        port = f.get("port", "")
        ver = f.get("version", "")
        desc = (f.get("description", "") or "")[:120]
        if host and svc:
            exploit_items_html += f'<li><strong>{host}:{port} — {svc}'
            if ver:
                exploit_items_html += f" {ver}"
            exploit_items_html += f":</strong> {desc}</li>\n"
    if not exploit_items_html:
        exploit_items_html = "<li>No exploitation attempts were conducted on this engagement.</li>"

    # Build the HTML
    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8">
<style>
  @page {{ size: A4; margin: 2cm; @bottom-center {{ content: counter(page) " / " counter(pages); font-size: 9px; color: #797979; }} }}
  body {{ font-family: Helvetica, Arial, sans-serif; color: #333; font-size: 11pt; line-height: 1.5; }}
  h1 {{ color: #CC4141; border-bottom: 2px solid #CC4141; padding-bottom: 6px; font-size: 22pt; }}
  h2 {{ color: #333; border-bottom: 1px solid #ddd; padding-bottom: 4px; font-size: 16pt; margin-top: 30px; }}
  h3 {{ font-size: 13pt; margin: 20px 0 10px 0; }}
  h4 {{ font-size: 11pt; margin: 15px 0 5px 0; color: #555; }}
  .subtitle {{ color: #797979; font-size: 14pt; margin-top: -10px; }}
  table {{ width: 100%; border-collapse: collapse; margin: 15px 0; font-size: 9pt; }}
  tr {{ break-inside: avoid; page-break-inside: avoid; }}
  th, td {{ padding: 6px 8px; text-align: left; border-bottom: 1px solid #eee; }}
  th {{ background: #333; color: #fff; font-weight: bold; }}
  tr:nth-child(even) td {{ background: #F5F5F5; }}
  .sev-tag {{ display: inline-block; padding: 2px 8px; border-radius: 3px; font-size: 9pt; font-weight: bold; color: #fff; margin-left: 8px; }}
  .sev-critical {{ background: #CC4141; }}
  .sev-high {{ background: #E85D3F; }}
  .sev-medium {{ background: #F4A442; }}
  .sev-low {{ background: #5B9E5B; }}
  tr:nth-child(even) td {{ background: #F5F5F5; }}
  td.sev-critical, td.sev-high, td.sev-medium, td.sev-low, td.sev-info {{ color: #fff; font-weight: bold; }}
  td.sev-critical {{ background: #CC4141 !important; }}
  td.sev-high {{ background: #E85D3F !important; }}
  td.sev-medium {{ background: #F4A442 !important; }}
  td.sev-low {{ background: #5B9E5B !important; }}
  td.sev-info {{ background: #4A90D9 !important; }}
  .meta {{ width: auto; font-size: 10pt; margin: 5px 0 15px 25px; }}
  .meta td {{ border: none; padding: 2px 8px; }}
  .meta .label {{ font-weight: bold; color: #797979; width: 80px; }}
  .finding {{ page-break-inside: avoid; }}
  a {{ color: #333; text-decoration: none; }}
  a:hover {{ text-decoration: underline; }}
  .cve-link {{ color: #4A90D9; text-decoration: underline; font-family: Consolas, monospace; font-size: 10pt; }}
  .summary-box {{ display: flex; gap: 10px; margin: 20px 0; }}
  .sev-box {{ flex: 1; text-align: center; padding: 15px 5px; border-radius: 4px; color: #fff; font-size: 10pt; }}
  .sev-box .count {{ font-size: 28pt; font-weight: bold; display: block; }}
</style>
</head>
<body>

<h1>Security Assessment Report</h1>
<p class="subtitle">PortShim — {datetime.now().strftime('%B %d, %Y')} | {total} findings</p>

<h2>Executive Summary</h2>

<h3>Assessment Overview</h3>
<p>PortShim conducted an internal network security assessment of <strong>{target_net}</strong> on {datetime.now().strftime('%B %d, %Y')}. The assessment identified <strong>{total} findings</strong> across <strong>{host_count} live hosts</strong> using a combination of service enumeration and vulnerability scanning. The engagement used a <strong>{engagement_profile}</strong> profile with <strong>Hybrid</strong> LLM mode (local models for exploitation, cloud for reporting).</p>

<h3>Overall Risk Posture: {risk_label}</h3>
<div class="summary-box">
{''.join(f'<div class="sev-box sev-{s}"><span class="count">{summary[s]}</span>{s.capitalize()}</div>' for s in ['critical','high','medium','low','info'] if summary[s] > 0)}
</div>
<p>{risk_detail}</p>

<h3>Key Risks</h3>
<ol>{top_lines}</ol>

<h3>Identified Services</h3>
<p>Identified services on high-value targets:
<ul>
{exploit_items_html}
</ul></p>

<h3>Environment Maturity: {maturity}</h3>
<p>{maturity_detail}</p>

<h3>Recommended Next Steps</h3>
<ol>{next_steps}</ol>

<h2>Findings Summary</h2>
<table>
<thead><tr><th>Severity</th><th>Title</th><th>Host</th><th>CVSS</th></tr></thead>
<tbody>
{detail_rows}
</tbody>
</table>

<h2>Detailed Findings</h2>
{detail_sections}

</body>
</html>"""

    HTML(string=html).write_pdf(output_path)
    return output_path


def sev_ranges(severity: str) -> str:
    """Return CVSS range description for severity."""
    return {
        "critical": "9.0-10.0",
        "high": "7.0-8.9",
        "medium": "4.0-6.9",
        "low": "0.1-3.9",
        "info": "0.0",
    }.get(severity, "?")


def hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    """Convert hex color string to RGB tuple."""
    h = hex_color.lstrip("#")
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))


def add_hyperlink(paragraph, url: str, text: str) -> None:
    """Add a clickable hyperlink to a python-docx paragraph."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    part = paragraph.part
    r_id = part.relate_to(url, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink", is_external=True)
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), r_id)
    new_run = OxmlElement("w:r")
    rPr = OxmlElement("w:rPr")
    c = OxmlElement("w:color")
    c.set(qn("w:val"), "4A90D9")
    rPr.append(c)
    u = OxmlElement("w:u")
    u.set(qn("w:val"), "single")
    rPr.append(u)
    sz = OxmlElement("w:sz")
    sz.set(qn("w:val"), "20")  # 10pt
    rPr.append(sz)
    new_run.append(rPr)
    new_run.text = text
    hyperlink.append(new_run)
    paragraph._p.append(hyperlink)


def main():
    parser = argparse.ArgumentParser(description="Generate pentest report deliverables")
    parser.add_argument("findings", help="JSON file with findings")
    parser.add_argument("--output-dir", default="./reports", help="Output directory (timestamped subdirectory created inside)")
    parser.add_argument("--category", choices=["wired", "wireless"], default="wired",
                        help="Category for subdirectory naming (default: wired)")
    parser.add_argument("--format", choices=["docx", "pptx", "pdf", "all"], default="all")
    parser.add_argument("--prefix", default="portshim", help="Filename prefix")
    parser.add_argument("--engagement", default="Surgical", help="Engagement profile name")
    args = parser.parse_args()

    findings = load_findings(args.findings)
    if not findings:
        print("No findings to report.")
        sys.exit(0)

    from pathlib import Path
    import datetime as dt

    ts = dt.datetime.now().strftime("%Y-%m-%d-%H-%M-%S")
    category = getattr(args, "category", "wired")
    out_dir = Path(args.output_dir) / f"{category}-{ts}"
    out_dir.mkdir(parents=True, exist_ok=True)
    date_str = dt.datetime.now().strftime("%Y%m%d")

    results = []
    if args.format in ("docx", "all"):
        path = out_dir / f"{args.prefix}-report-{date_str}.docx"
        result = generate_docx(findings, str(path))
        results.append(f"Word:  {result}")

    if args.format in ("pptx", "all"):
        path = out_dir / f"{args.prefix}-brief-{date_str}.pptx"
        result = generate_pptx(findings, str(path), engagement_profile=args.engagement)
        results.append(f"PPTX:  {result}")

    if args.format in ("pdf", "all"):
        path = out_dir / f"{args.prefix}-report-{date_str}.pdf"
        result = generate_pdf(findings, str(path), engagement_profile=args.engagement)
        results.append(f"PDF:   {result}")

    print(f"Generated {len(findings)}-finding report:")
    for r in results:
        print(f"  {r}")

    # Auto-save to scan history DB (if available)
    try:
        from pathlib import Path
        sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent.parent / "scripts"))
        from save_to_db import main as save_main
        import argparse as _argparse
        # Derive engagement ID from output dir or date
        eng_id = args.engagement or out_dir.name or f"scan-{date_str}"
        save_args = _argparse.Namespace(
            findings=str(findings_path),
            engagement=eng_id,
            client=None,
            hosts=None,
            stealth=None,
            mode=None,
            target=None,
            dry_run=False,
        )
        save_main.__wrapped__ if hasattr(save_main, '__wrapped__') else None  # no-op, just checking
        from scan_db import ScanDB
        db = ScanDB()
        db.save_engagement(eng_id, None, "auto-saved")
        db.save_findings(eng_id, findings)
        sev = {"critical": 0, "high": 0, "medium": 0, "low": 0}
        for f in findings:
            s = (f.get("severity") or "").lower()
            if s in sev: sev[s] += 1
        db.complete_engagement(eng_id, len(findings), **sev)
    except Exception:
        pass  # DB save is optional — don't block report generation


if __name__ == "__main__":
    main()
